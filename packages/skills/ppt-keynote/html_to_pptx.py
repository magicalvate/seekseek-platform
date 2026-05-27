"""
Convert ppt-keynote HTML deck → editable .pptx (WPS / PowerPoint compatible).

Handles the full content from each slide type:
  - cover       : kicker + big title + subtitle
  - stat-row    : kicker + title + 3-column stat cards (num + label)
  - point-list  : kicker + title + up to 3 icon/title/desc rows
  - big-quote   : large quote text + source
  - thanks      : oversized closing word + subtitle

Usage:
    python3 html_to_pptx.py input.html output.pptx
"""

import sys
import re
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Canvas ────────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)

# Rails
MX          = Inches(0.85)          # left/right margin
CONTENT_W   = SLIDE_W - MX * 2
MARGIN_TOP  = Inches(0.7)
FOOTER_TOP  = Inches(6.9)

# ── Palette ───────────────────────────────────────────────────────────────────
WHITE   = RGBColor(0xF9, 0xFA, 0xFB)
LIGHT   = RGBColor(0xF3, 0xF4, 0xF6)
MUTED   = RGBColor(0x6B, 0x72, 0x80)
DIM     = RGBColor(0x37, 0x41, 0x51)
ACCENT  = RGBColor(0xA7, 0x8B, 0xFA)   # purple
GREEN   = RGBColor(0x34, 0xD3, 0x99)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)

# Per-slide background gradient approximations (start color used as solid)
BG_PALETTE = [
    RGBColor(0x17, 0x15, 0x35),   # 1 cover     – deep purple
    RGBColor(0x0D, 0x14, 0x28),   # 2 dark blue
    RGBColor(0x0A, 0x12, 0x28),   # 3 navy
    RGBColor(0x11, 0x11, 0x11),   # 4 near-black
    RGBColor(0x09, 0x14, 0x0E),   # 5 dark green
    RGBColor(0x18, 0x09, 0x09),   # 6 dark red
    RGBColor(0x0A, 0x0A, 0x1F),   # 7 indigo
    RGBColor(0x0D, 0x0D, 0x0D),   # 8 black
]

ACCENT_CYCLE = [ACCENT, GREEN, AMBER]   # rotate through stat cards


# ── HTML helpers ──────────────────────────────────────────────────────────────

def _strip(html: str) -> str:
    """Remove tags, normalise whitespace, decode entities."""
    text = re.sub(r'<br\s*/?>', '\n', html or '')
    text = re.sub(r'<[^>]+>', '', text)
    for ent, ch in [('&amp;', '&'), ('&lt;', '<'), ('&gt;', '>'), ('&nbsp;', ' '), ('&#x27;', "'")]:
        text = text.replace(ent, ch)
    return re.sub(r'[ \t]+', ' ', text).strip()


def _find(pattern: str, html: str) -> str:
    m = re.search(pattern, html, re.DOTALL)
    return _strip(m.group(1)) if m else ''


def parse_slides(html: str) -> list[dict]:
    sections = re.findall(
        r'<section[^>]+class="[^"]*slide[^"]*"[^>]*>(.*?)</section>',
        html, re.DOTALL
    )
    slides = []
    for raw in sections:
        kicker   = _find(r'class="kicker"[^>]*>(.*?)</div>', raw)
        title    = _strip(re.sub(r'<span[^>]*>.*?</span>', '', _find(r'<h[12][^>]*>(.*?)</h[12]>', raw)) )
        subtitle = _find(r'class="subtitle"[^>]*>(.*?)</p>', raw)

        # stat cards
        stats = []
        for card in re.findall(r'<div class="stat"[^>]*>(.*?)</div>\s*</div>', raw, re.DOTALL):
            num   = _find(r'class="stat-num"[^>]*>(.*?)</div>', card)
            label = _find(r'class="stat-label"[^>]*>(.*?)</div>', card)
            if num or label:
                stats.append({'num': num, 'label': label})

        # point rows
        points = []
        for pt in re.findall(r'<div class="point"[^>]*>(.*?)</div>\s*</div>', raw, re.DOTALL):
            pt_title = _find(r'class="point-title"[^>]*>(.*?)</div>', pt)
            pt_desc  = _find(r'class="point-desc"[^>]*>(.*?)</div>', pt)
            if pt_title:
                points.append({'title': pt_title, 'desc': pt_desc})

        # big quote
        big_quote = _find(r'class="big-quote"[^>]*>(.*?)</div>', raw)
        quote_src = _find(r'class="quote-src"[^>]*>(.*?)</div>', raw)

        # thanks
        thanks     = _find(r'class="thanks"[^>]*>(.*?)</div>', raw)
        thanks_sub = _find(r'class="thanks-sub"[^>]*>(.*?)</p>', raw)

        # classify slide type
        if thanks:
            kind = 'thanks'
        elif big_quote:
            kind = 'quote'
        elif stats:
            kind = 'stats'
        elif points:
            kind = 'points'
        else:
            kind = 'cover'

        slides.append({
            'kind': kind,
            'kicker': kicker, 'title': title, 'subtitle': subtitle,
            'stats': stats, 'points': points,
            'big_quote': big_quote, 'quote_src': quote_src,
            'thanks': thanks, 'thanks_sub': thanks_sub,
        })
    return slides


# ── python-pptx helpers ───────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txb(slide, text: str, l, t, w, h,
        size: int, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True) -> None:
    if not text:
        return
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size     = Pt(size)
        run.font.bold     = bold
        run.font.italic   = italic
        run.font.color.rgb = color
        run.font.name     = 'Arial'


def card_rect(slide, l, t, w, h, fill=RGBColor(0x18,0x18,0x24), border=RGBColor(0x2d,0x2d,0x44)):
    """Rounded rectangle card background."""
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb      = border
    shape.line.width          = Pt(0.5)
    return shape


def page_num(slide, cur, total):
    txb(slide, f'{cur} / {total}',
        SLIDE_W - Inches(1.4), FOOTER_TOP,
        Inches(1.2), Inches(0.4),
        size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ── Slide builders ────────────────────────────────────────────────────────────

def build_cover(slide, d, idx, total):
    y = Inches(1.6)
    if d['kicker']:
        txb(slide, d['kicker'].upper(), MX, y, CONTENT_W, Inches(0.4),
            size=10, color=MUTED)
        y += Inches(0.5)
    if d['title']:
        txb(slide, d['title'], MX, y, CONTENT_W, Inches(2.4),
            size=66, bold=True, color=WHITE)
        y += Inches(2.5)
    if d['subtitle']:
        txb(slide, d['subtitle'], MX, y, Inches(8), Inches(1.2),
            size=22, color=MUTED)
    page_num(slide, idx, total)


def build_stats(slide, d, idx, total):
    y = MARGIN_TOP
    if d['kicker']:
        txb(slide, d['kicker'].upper(), MX, y, CONTENT_W, Inches(0.4), size=10, color=MUTED)
        y += Inches(0.48)
    if d['title']:
        txb(slide, d['title'], MX, y, CONTENT_W, Inches(1.3), size=46, bold=True, color=WHITE)
        y += Inches(1.45)

    cards = d['stats'][:3]
    if not cards:
        page_num(slide, idx, total); return

    gap    = Inches(0.18)
    card_w = (CONTENT_W - gap * (len(cards) - 1)) / len(cards)
    card_h = Inches(2.0)
    x = MX

    for i, stat in enumerate(cards):
        card_rect(slide, x, y, card_w, card_h)
        accent = ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        # number
        txb(slide, stat['num'], x + Inches(0.22), y + Inches(0.22), card_w - Inches(0.44), Inches(0.9),
            size=44, bold=True, color=accent)
        # label
        txb(slide, stat['label'], x + Inches(0.22), y + Inches(1.1), card_w - Inches(0.44), Inches(0.75),
            size=14, color=MUTED)
        x += card_w + gap

    page_num(slide, idx, total)


def build_points(slide, d, idx, total):
    y = MARGIN_TOP
    if d['kicker']:
        txb(slide, d['kicker'].upper(), MX, y, CONTENT_W, Inches(0.4), size=10, color=MUTED)
        y += Inches(0.48)
    if d['title']:
        txb(slide, d['title'], MX, y, CONTENT_W, Inches(1.1), size=44, bold=True, color=WHITE)
        y += Inches(1.25)

    pts   = d['points'][:3]
    gap   = Inches(0.14)
    row_h = min(Inches(1.5), (Inches(6.8) - y - gap * (len(pts)-1)) / len(pts))

    for pt in pts:
        card_rect(slide, MX, y, CONTENT_W, row_h)
        # title
        txb(slide, pt['title'],
            MX + Inches(0.25), y + Inches(0.18),
            CONTENT_W - Inches(0.5), Inches(0.5),
            size=17, bold=True, color=LIGHT)
        # desc
        if pt['desc']:
            txb(slide, pt['desc'],
                MX + Inches(0.25), y + Inches(0.62),
                CONTENT_W - Inches(0.5), row_h - Inches(0.72),
                size=14, color=MUTED)
        y += row_h + gap

    page_num(slide, idx, total)


def build_quote(slide, d, idx, total):
    y = Inches(1.6)
    if d['big_quote']:
        txb(slide, d['big_quote'], MX, y, CONTENT_W, Inches(3.2),
            size=38, bold=True, color=WHITE)
        y += Inches(3.4)
    if d['quote_src']:
        txb(slide, d['quote_src'], MX, y, CONTENT_W, Inches(0.5),
            size=15, color=DIM)
    page_num(slide, idx, total)


def build_thanks(slide, d, idx, total):
    y = Inches(1.8)
    if d['thanks']:
        txb(slide, d['thanks'], MX, y, CONTENT_W, Inches(2.2),
            size=84, bold=True, color=WHITE)
        y += Inches(2.4)
    if d['thanks_sub']:
        txb(slide, d['thanks_sub'], MX, y, CONTENT_W, Inches(0.6),
            size=18, color=DIM)
    page_num(slide, idx, total)


BUILDERS = {
    'cover':  build_cover,
    'stats':  build_stats,
    'points': build_points,
    'quote':  build_quote,
    'thanks': build_thanks,
}


# ── Main ──────────────────────────────────────────────────────────────────────

def build_pptx(slides_data: list[dict], out_path: Path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = len(slides_data)

    for i, d in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        set_bg(slide, BG_PALETTE[i % len(BG_PALETTE)])
        builder = BUILDERS.get(d['kind'], build_cover)
        builder(slide, d, i + 1, total)

    prs.save(str(out_path))
    print(f'✅  {out_path}  ({total} 张)')


if __name__ == '__main__':
    src = Path(sys.argv[1] if len(sys.argv) > 1 else 'verify-output.html')
    dst = Path(sys.argv[2] if len(sys.argv) > 2 else src.with_suffix('.pptx'))

    html = src.read_text(encoding='utf-8')
    data = parse_slides(html)

    if not data:
        print('❌  未识别到 <section class="slide">'); sys.exit(1)

    print(f'识别到 {len(data)} 张幻灯片：')
    for j, s in enumerate(data, 1):
        print(f'  {j}. [{s["kind"]:7}] {(s["title"] or s["thanks"] or s["big_quote"] or "—")[:45]}')

    build_pptx(data, dst)
